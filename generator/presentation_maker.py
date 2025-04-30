import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageFilter
from io import BytesIO
import re
import os
from transformers import pipeline
from generator.summarizer import smart_conclusion_human_style


try:
    summarizer = pipeline("summarization", model="IlyaGusev/rut5_base_sum_gazeta")
except:
    summarizer = None

def summarize_to_bullets(text, max_sentences=4):
    if summarizer is None:
        return []

    try:
        result = summarizer(text, max_length=120, min_length=40, do_sample=False)
        summary = result[0]['summary_text']
        bullets = re.split(r'(?<=[.!?])\s+', summary.strip())
        return bullets[:max_sentences]
    except:
        return []

def split_text_on_slides(text, max_characters=800):
    """
    Разбивает текст на несколько частей, если он превышает заданный лимит по символам.
    """
    slides = []
    while len(text) > max_characters:
        split_point = text.rfind(" ", 0, max_characters)  # Находим последний пробел в пределах лимита
        slides.append(text[:split_point].strip())  # Добавляем часть текста
        text = text[split_point:].strip()  # Оставшийся текст
    slides.append(text)  # Добавляем последнюю часть текста
    return slides

def generate_presentation(title, intro, sections, conclusion, image_urls):
    prs = Presentation()

    # Слайд с заголовком
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Презентация по теме"

    # Введение: разделяем картинку и текст
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Пустой слайд
    if image_urls:
        try:
            response = requests.get(image_urls[0])
            if response.status_code == 200 and "image" in response.headers.get("Content-Type", ""):
                img = Image.open(BytesIO(response.content))
                blurred = img.filter(ImageFilter.GaussianBlur(radius=15))  # 🌫 Размытие

                blurred_path = "temp_blurred.jpg"
                blurred.save(blurred_path)

                slide.shapes.add_picture(
                    blurred_path, Inches(0), Inches(0),
                    width=prs.slide_width, height=prs.slide_height
                )
        except Exception as e:
            print(f"⚠ Ошибка загрузки или обработки фона: {e}")
    

    title_box = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(9), Inches(1.5))
    tf_title = title_box.text_frame
    p = tf_title.paragraphs[0]
    p.text = "Введение"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # Белый
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Текст (ввод)
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(5))
    tf_intro = intro_box.text_frame
    tf_intro.word_wrap = True
    tf_intro.text_anchor = MSO_ANCHOR.TOP
    p = tf_intro.add_paragraph()
    p.text = intro
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.line_spacing = 1.3

    # Основная часть
    for i, (sec_title, sec_text) in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Слайд с пустыми полями
        slide.shapes.title.text = sec_title
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 200)  # Светло-желтый фон

        # Текст
        tf = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(5), Inches(8))
        text_frame = tf.text_frame
        text_frame.word_wrap = True
        text_frame.text_anchor = MSO_ANCHOR.TOP

        bullets = summarize_to_bullets(sec_text)

        if not bullets:
            bullets = ["Информация временно недоступна."]

        for sentence in bullets:
            p = text_frame.add_paragraph()
            p.text = sentence.strip()
            p.font.size = Pt(18)
            p.level = 0
            p.space_after = Pt(6)
            p.line_spacing = 1.2



        # Картинка
        if i < len(image_urls):
            try:
                response = requests.get(image_urls[i])
                if response.status_code == 200 and "image" in response.headers.get("Content-Type", ""):
                    img_stream = BytesIO(response.content)
                    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.5), height=Inches(5), width=Inches(4))
            except:
                pass

    # Заключение
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Слайд с пустыми полями
    slide.shapes.title.text = "Заключение"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Розовый фон

    # Текст
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(5))
    text_frame = tf.text_frame
    text_frame.word_wrap = True  # Включаем автоматический перенос текста
    p = text_frame.add_paragraph()
    p.text = conclusion
    p.font.size = Pt(18)
    text_frame.text_anchor = MSO_ANCHOR.TOP  # Чтобы текст не сжимался внизу

    filename = f"{title}_presentation.pptx"
    prs.save(filename)
    if os.path.exists("temp_blurred.jpg"):
        os.remove("temp_blurred.jpg")