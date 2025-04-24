import wikipedia
import re
import json
import requests
import urllib.parse
from io import BytesIO
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
import dash
from dash import dcc, html, Input, Output, State
import dash_bootstrap_components as dbc
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR


wikipedia.set_lang("ru")
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
app = dash.Dash(__name__, external_stylesheets=[dbc.themes.BOOTSTRAP])
server = app.server

TOPIC_TRANSLATIONS = {
    "горы": "mountains", "космос": "space", "море": "sea", "океан": "ocean",
    "животные": "animals", "человек": "human", "природа": "nature", "город": "city",
    "архитектура": "architecture", "история": "history", "техника": "technology",
    "наука": "science", "музыка": "music", "спорт": "sport", "еда": "food"
}

def get_clean_article(title):
    try:
        page = wikipedia.page(title)
        text = page.content
        text = re.split(r"==\s*См\. также\s*==", text)[0]
        text = re.sub(r"==+\s*(.*?)\s*==+", r"§\1§", text)
        return text.strip()
    except:
        return None

def split_into_sections(text, max_sections):
    sections = []
    parts = re.split(r'§(.*?)§', text)
    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        content = parts[i + 1].strip()
        sections.append((title, content))
    return sections[:max_sections]

def chunk_text_to_bullets(text, max_lines=4):
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    return [s.strip() for s in sentences if len(s.strip()) > 20][:max_lines]

def fetch_image_urls_bing(query, count):
    search_term = TOPIC_TRANSLATIONS.get(query.lower(), query)
    headers = {"User-Agent": "Mozilla/5.0"}
    url = f"https://www.bing.com/images/search?q={urllib.parse.quote(search_term)}"

    response = requests.get(url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    image_tags = soup.find_all("a", class_="iusc")

    urls = []
    for tag in image_tags:
        try:
            m_json = json.loads(tag.get("m"))
            img_url = m_json["murl"]
            if "ytimg.com" in img_url or "youtube" in img_url:
                continue
            if img_url.endswith((".jpg", ".jpeg", ".png")):
                urls.append(img_url)
            if len(urls) >= count:
                break
        except:
            continue
    return urls

try:
    summarizer = pipeline(
        "summarization",
        model="IlyaGusev/rut5_base_sum_gazeta",
        tokenizer="IlyaGusev/rut5_base_sum_gazeta"
    )
except:
    # Запасной вариант если модель не загрузится
    summarizer = None

def smart_conclusion(title, sections):
    print("🧠 Генерируем осмысленное заключение...")
    
    # Варианты заключений на случай ошибок
    FALLBACK_CONCLUSIONS = [
        "На основе представленных материалов можно сделать вывод о важности и актуальности данной темы.",
        "Проведенный анализ позволяет утверждать, что эта тема заслуживает особого внимания.",
        "Таким образом, рассмотренная тема представляет значительный научный и практический интерес.",
        "Изучение данной темы позволяет лучше понять ключевые аспекты этой области знаний."
    ]
    
    # Собираем текст из всех разделов
    combined = "\n".join([f"{sec[0]}: {sec[1]}" for sec in sections if len(sec[1].strip()) > 30])
    
    # Если текста совсем мало, возвращаем запасной вариант
    if len(combined) < 200:
        return FALLBACK_CONCLUSIONS[0]
    
    # Нормализуем текст
    clean_text = ' '.join(combined.replace("\n", " ").split())
    
    # Если модель не загрузилась, используем запасной вариант
    if summarizer is None:
        return FALLBACK_CONCLUSIONS[1]
    
    try:
        # Генерируем ключевые пункты
        key_points = summarizer(
            clean_text,
            max_length=150,
            min_length=50,
            do_sample=False,
            truncation=True
        )[0]['summary_text']
        
        # На основе ключевых пунктов генерируем заключение
        conclusion = summarizer(
            f"Ключевые пункты: {key_points}. Напиши заключение научного доклада.",
            max_length=120,
            min_length=60,
            do_sample=False,
            truncation=True
        )[0]['summary_text']
        
        # Проверяем результат
        if len(conclusion) > 20 and "." in conclusion:
            return conclusion
        else:
            return FALLBACK_CONCLUSIONS[2]
            
    except Exception as e:
        print(f"⚠ Ошибка генерации: {str(e)}")
        return FALLBACK_CONCLUSIONS[3]

def generate_report(title, intro, sections, conclusion):
    doc = Document()
    doc.add_heading(title, 0)

    doc.add_heading("Введение", level=1)
    doc.add_paragraph(intro)

    doc.add_heading("Основная часть", level=1)
    for sec_title, sec_text in sections:
        doc.add_heading(sec_title, level=2)
        doc.add_paragraph(sec_text)

    doc.add_heading("Заключение", level=1)
    final_thought = smart_conclusion(title, sections)
    if len(final_thought) < 30:  # Если заключение слишком короткое
        final_thought = "Таким образом, рассмотренная тема представляет собой важный аспект для изучения."

    doc.add_paragraph(final_thought)

    filename = f"{title}_report.docx"
    doc.save(filename)


def split_text_on_slides(text, max_characters=800):
    """
    Разбивает текст на несколько частей, если он превышает заданный лимит по символам.
    """
    slides = []
    while len(text) > max_characters:
        split_point = text.rfind(" ", 0, max_characters)  # Находим последний пробел в пределах лимита
        slides.append(text[:split_point].strip())  # Добавляем часть текста
        text = text[split_point:].strip()  # Оставшийся текст
    slides.append(text)  # Добавляем последнюю часть текста
    return slides

def generate_presentation(title, intro, sections, conclusion, image_urls):
    prs = Presentation()

    # Слайд с заголовком
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Презентация по теме"

    # Введение: разделяем картинку и текст
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Используем layout с пустыми полями
    slide.shapes.title.text = "Введение"
    slide.background.fill.solid()  # Устанавливаем фон слайда
    slide.background.fill.fore_color.rgb = RGBColor(200, 220, 255)  # Синий оттенок

    # Текст
    tf = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(5))
    text_frame = tf.text_frame
    text_frame.word_wrap = True  # Включаем автоматический перенос текста
    p = text_frame.add_paragraph()
    p.text = intro
    p.font.size = Pt(18)
    text_frame.text_anchor = MSO_ANCHOR.TOP  # Чтобы текст не сжимался внизу

    # Картинка
    if image_urls:
        try:
            response = requests.get(image_urls[0])
            if response.status_code == 200 and "image" in response.headers.get("Content-Type", ""):
                img_stream = BytesIO(response.content)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(4.5), height=Inches(3), width=Inches(3))
        except:
            pass

    # Основная часть
    for i, (sec_title, sec_text) in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Слайд с пустыми полями
        slide.shapes.title.text = sec_title
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 200)  # Светло-желтый фон

        # Текст
        tf = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(5))
        text_frame = tf.text_frame
        text_frame.word_wrap = True  # Включаем автоматический перенос текста
        p = text_frame.add_paragraph()
        p.text = sec_text
        p.font.size = Pt(18)
        text_frame.text_anchor = MSO_ANCHOR.TOP  # Чтобы текст не сжимался внизу

        # Картинка
        if i < len(image_urls):
            try:
                response = requests.get(image_urls[i])
                if response.status_code == 200 and "image" in response.headers.get("Content-Type", ""):
                    img_stream = BytesIO(response.content)
                    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(4.5), height=Inches(3), width=Inches(3))
            except:
                pass

    # Заключение
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Слайд с пустыми полями
    slide.shapes.title.text = "Заключение"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Розовый фон

    # Текст
    tf = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(5))
    text_frame = tf.text_frame
    text_frame.word_wrap = True  # Включаем автоматический перенос текста
    p = text_frame.add_paragraph()
    p.text = conclusion
    p.font.size = Pt(18)
    text_frame.text_anchor = MSO_ANCHOR.TOP  # Чтобы текст не сжимался внизу

    filename = f"{title}_presentation.pptx"
    prs.save(filename)

app.layout = dbc.Container([
    html.H2("Автоматический генератор доклада и презентации"),
    dbc.Row([
        dbc.Col([
            html.Label("Тема"),
            dcc.Input(id="topic", type="text", value="Космос", style={"width": "100%"}),
            html.Br(), html.Br(),
            html.Label("Что сгенерировать:"),
            dcc.RadioItems(
                id="mode", options=[
                    {"label": "Доклад", "value": "report"},
                    {"label": "Презентация", "value": "presentation"},
                    {"label": "Оба", "value": "both"},
                ],
                value="both",
                labelStyle={'display': 'block'}
            ),
            html.Label("Детализация текста:"),
            dcc.Dropdown(
                options=[
                    {"label": "Краткий", "value": "краткий"},
                    {"label": "Средний", "value": "средний"},
                    {"label": "Подробный", "value": "подробный"}
                ],
                value="средний",
                id="detail"
            ),
            html.Label("Слайдов (если презентация):"),
            dcc.Input(id="slides", type="number", value=8, min=5, max=15),
            html.Label("Картинок в презентации:"),
            dcc.Input(id="images", type="number", value=4, min=0, max=12),
            html.Br(), html.Br(),
            dbc.Button("Сгенерировать", id="generate", color="primary"),
        ], width=6)
    ]),
    html.Br(),
    html.Div(id="output")
])

@app.callback(
    Output("output", "children"),
    Input("generate", "n_clicks"),
    State("topic", "value"),
    State("mode", "value"),
    State("detail", "value"),
    State("slides", "value"),
    State("images", "value"),
)
def run_generator(n, topic, mode, detail, slides, images):
    if not n:
        return ""
    text = get_clean_article(topic)
    if not text:
        return "❌ Не удалось получить статью с Википедии."

    intro = text[:400]
    conclusion = text[-400:]
    sections = split_into_sections(text, slides - 3)

    if mode in ["report", "both"]:
        generate_report(topic, intro, sections, conclusion)
    if mode in ["presentation", "both"]:
        image_urls = fetch_image_urls_bing(topic, images)
        generate_presentation(topic, intro, sections, conclusion, image_urls)

    return f"✅ Готово! Файлы сохранены: {topic}_report.docx и/или {topic}_presentation.pptx"

if __name__ == "__main__":
    app.run_server(debug=True)